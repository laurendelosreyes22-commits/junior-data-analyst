"""
Generate Financial Keyword Intelligence presentation as .pptx
Run: python docs/build_slides.py
Output: docs/slides.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── Colors ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0d, 0x0d, 0x1a)
WHITE    = RGBColor(0xff, 0xff, 0xff)
INDIGO   = RGBColor(0x63, 0x66, 0xf1)
VIOLET   = RGBColor(0x8b, 0x5c, 0xf6)
GREEN    = RGBColor(0x10, 0xb9, 0x81)
SUBTLE   = RGBColor(0x99, 0x9a, 0xbb)
TAG      = RGBColor(0x55, 0x55, 0x88)
CALLOUT_I_BG = RGBColor(0x1a, 0x1a, 0x35)
CALLOUT_V_BG = RGBColor(0x1a, 0x14, 0x35)
CALLOUT_G_BG = RGBColor(0x0a, 0x2a, 0x20)
SCORE    = RGBColor(0xa7, 0x8b, 0xfa)
DIM      = RGBColor(0x44, 0x44, 0x66)

# ── Helpers ──────────────────────────────────────────────────────────────────
W  = Inches(13.33)   # slide width
H  = Inches(7.5)     # slide height
M  = Inches(1.1)     # left margin
MT = Inches(0.75)    # top margin

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank layout


def new_slide(accent_color=None):
    """Add a blank dark slide, optionally with a left accent bar."""
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG

    if accent_color:
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0),
            Inches(0.12), H
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()

    return slide


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_tag(slide, text):
    add_text(slide, text.upper(),
             M, MT, Inches(10), Inches(0.4),
             size=10, bold=True, color=TAG)


def add_h1(slide, text, y=Inches(1.4)):
    add_text(slide, text,
             M, y, Inches(11), Inches(2.5),
             size=52, bold=True, color=WHITE)


def add_h2(slide, text, y=Inches(1.4)):
    add_text(slide, text,
             M, y, Inches(11), Inches(2.2),
             size=36, bold=True, color=WHITE)


def add_sub(slide, text, y):
    add_text(slide, text,
             M, y, Inches(11), Inches(1),
             size=18, color=SUBTLE)


def add_callout(slide, text, y, bg_color, border_color):
    box = slide.shapes.add_shape(
        1, M, y, Inches(11), Inches(0.9)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = border_color
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"


def add_bar(slide, label, value, max_val, y, color):
    """Draw a labeled bar row."""
    lx = M
    add_text(slide, label, lx, y, Inches(1.8), Inches(0.45),
             size=13, color=SUBTLE, align=PP_ALIGN.RIGHT)

    track_x = M + Inches(1.9)
    track_w = Inches(6.5)

    # track background
    track = slide.shapes.add_shape(1, track_x, y, track_w, Inches(0.38))
    track.fill.solid()
    track.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x33)
    track.line.fill.background()

    # fill bar
    fill_w = int(track_w * value / max_val)
    fill = slide.shapes.add_shape(1, track_x, y, fill_w, Inches(0.38))
    fill.fill.solid()
    fill.fill.fore_color.rgb = color
    fill.line.fill.background()

    # value label
    add_text(slide, str(value),
             track_x + fill_w + Inches(0.1), y, Inches(0.5), Inches(0.38),
             size=13, bold=True, color=WHITE)


def add_slide_number(slide, n):
    add_text(slide, str(n),
             W - Inches(0.8), H - Inches(0.55), Inches(0.6), Inches(0.4),
             size=12, color=DIM, align=PP_ALIGN.RIGHT)


def add_appear_animations(slide, shape_groups):
    """Add sequential click-to-appear animations.
    shape_groups: list of lists of shapes — each inner list appears together on one click."""
    pns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ids = [1]

    def nid():
        v = ids[0]; ids[0] += 1; return str(v)

    def sub(parent, tag, **attrs):
        el = etree.SubElement(parent, '{%s}%s' % (pns, tag))
        for k, v in attrs.items():
            el.set(k, str(v))
        return el

    timing = etree.Element('{%s}timing' % pns)
    tnLst = sub(timing, 'tnLst')
    par0 = sub(tnLst, 'par')
    cTn0 = sub(par0, 'cTn', id=nid(), dur='indefinite',
               restart='whenNotActive', nodeType='tmRoot')
    ch0 = sub(cTn0, 'childTnLst')

    seq = sub(ch0, 'seq', concurrent='1', nextAc='seek')
    cTnS = sub(seq, 'cTn', id=nid(), dur='indefinite', nodeType='mainSeq')
    chS = sub(cTnS, 'childTnLst')

    bldLst = etree.Element('{%s}bldLst' % pns)

    for grp_idx, shapes in enumerate(shape_groups):
        par1 = sub(chS, 'par')
        cTn1 = sub(par1, 'cTn', id=nid(), fill='hold')
        sc1 = sub(cTn1, 'stCondLst')
        sub(sc1, 'cond', evt='onBegin', delay='indefinite')
        ch1 = sub(cTn1, 'childTnLst')

        par2 = sub(ch1, 'par')
        cTn2 = sub(par2, 'cTn', id=nid(), fill='hold')
        sc2 = sub(cTn2, 'stCondLst')
        sub(sc2, 'cond', delay='0')
        ch2 = sub(cTn2, 'childTnLst')

        for shape in shapes:
            sp_id = str(shape.shape_id)
            par3 = sub(ch2, 'par')
            cTn3 = sub(par3, 'cTn', id=nid(),
                       presetID='1', presetClass='entr', presetSubtype='0',
                       fill='hold', grpId=str(grp_idx), nodeType='clickEffect')
            sc3 = sub(cTn3, 'stCondLst')
            sub(sc3, 'cond', delay='0')
            ch3 = sub(cTn3, 'childTnLst')

            set_el = sub(ch3, 'set')
            cBhvr = sub(set_el, 'cBhvr')
            sub(cBhvr, 'cTn', id=nid(), dur='1', fill='hold')
            tgtEl = sub(cBhvr, 'tgtEl')
            sub(tgtEl, 'spTgt', spid=sp_id)
            attrNl = sub(cBhvr, 'attrNameLst')
            attrN = sub(attrNl, 'attrName')
            attrN.text = 'style.visibility'
            to_el = sub(set_el, 'to')
            sub(to_el, 'strVal', val='visible')

            sub(bldLst, 'bldP', spid=sp_id, grpId=str(grp_idx),
                uiExpand='1', build='p')

    timing.append(bldLst)
    slide._element.append(timing)


# ── S1: HOOK ────────────────────────────────────────────────────────────────
s1 = new_slide()
add_tag(s1, "Hook")
add_h1(s1, "Right now, someone in Mississippi\nis searching for a payday loan.", y=Inches(1.5))
add_sub(s1, "Does your ad budget know that?", y=Inches(4.3))
add_slide_number(s1, 1)

# ── S2: THE PROBLEM ──────────────────────────────────────────────────────────
s2 = new_slide(accent_color=INDIGO)
add_tag(s2, "The Problem")
add_h2(s2, "EPCVIP sells leads to lenders.\nLeads are worth more in high-demand markets.")
add_callout(s2,
            "So which markets have the highest demand — and why?",
            Inches(4.0), CALLOUT_I_BG, INDIGO)
add_slide_number(s2, 2)

# ── S3: WHAT I BUILT ─────────────────────────────────────────────────────────
s3 = new_slide()
add_tag(s3, "What I Built")
add_h2(s3, "A pipeline that tracks financial keyword demand\nacross every US state — automatically.")

# Pipeline pills — track shapes for click-to-appear animations
pills = ["pytrends API", "→", "GitHub Actions", "→", "Snowflake", "→", "dbt", "→", "Streamlit Dashboard"]
px = M
py = Inches(4.1)
pill_groups = []   # each entry = [shape, ...] that appears on one click
pending_arrow = None

for pill in pills:
    is_arrow = pill == "→"
    w = Inches(0.35) if is_arrow else Inches(1.55)
    if not is_arrow:
        box = s3.shapes.add_shape(1, px, py, w, Inches(0.42))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x33)
        box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = pill
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = SUBTLE
        run.font.name = "Calibri"
        # arrow (if any) appears together with the pill it leads into
        pill_groups.append([pending_arrow, box] if pending_arrow else [box])
        pending_arrow = None
    else:
        pending_arrow = add_text(s3, pill, px + Inches(0.02), py, w, Inches(0.42),
                                 size=13, color=DIM, align=PP_ALIGN.CENTER)
    px += w + Inches(0.1)

add_sub(s3, "Runs daily. No manual work. Always fresh.", y=Inches(4.9))
add_slide_number(s3, 3)

# Animate: each pill (+ its preceding arrow) appears on a separate click
add_appear_animations(s3, pill_groups)

# ── S4: DESCRIPTIVE ──────────────────────────────────────────────────────────
s4 = new_slide(accent_color=INDIGO)
add_tag(s4, "Descriptive — What happened?")
add_h2(s4, "Personal loans lead nationally.\nBut the map tells a different story.")

bar_data = [
    ("Personal Loans",    52, INDIGO),
    ("Cash Advance",      50, VIOLET),
    ("Credit Cards",      42, RGBColor(0xa7, 0x8b, 0xfa)),
    ("Payday Loans",      27, VIOLET),
    ("Installment Loans", 19, INDIGO),
]
for i, (label, val, color) in enumerate(bar_data):
    add_bar(s4, label, val, 55, Inches(3.6 + i * 0.58), color)

# Dashboard map screenshot
import os
map_img = os.path.join(os.path.dirname(__file__), "dashboard-map.png")
s4.shapes.add_picture(map_img, Inches(8.9), Inches(3.3), Inches(4.0), Inches(3.5))

add_slide_number(s4, 4)

# ── S5: DIAGNOSTIC ───────────────────────────────────────────────────────────
s5 = new_slide(accent_color=VIOLET)
add_tag(s5, "Diagnostic — Why did it happen?")
add_h2(s5, "The South is underbanked.\nThat's where the demand is.")

# Table header
header_y = Inches(3.55)
for col, (label, cx, cw) in enumerate([
    ("Category",  M,              Inches(3.0)),
    ("Top State",  M + Inches(3.2), Inches(2.2)),
    ("Score",      M + Inches(5.6), Inches(1.0)),
]):
    add_text(s5, label.upper(), cx, header_y, cw, Inches(0.35),
             size=10, bold=True, color=TAG)

# Divider line
div = s5.shapes.add_shape(1, M, Inches(3.9), Inches(6.8), Inches(0.02))
div.fill.solid()
div.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x44)
div.line.fill.background()

rows = [
    ("Short-term credit", "Mississippi", "100"),
    ("Revolving credit",  "Wyoming",     "100"),
    ("Consumer lending",  "Mississippi", "84.5"),
]
for i, (cat, state, score) in enumerate(rows):
    ry = Inches(4.0 + i * 0.62)
    add_text(s5, cat,   M,              ry, Inches(3.0), Inches(0.5), size=16, color=WHITE)
    add_text(s5, state, M + Inches(3.2), ry, Inches(2.2), Inches(0.5), size=16, color=WHITE)
    add_text(s5, score, M + Inches(5.6), ry, Inches(1.0), Inches(0.5), size=16, bold=True, color=SCORE)

# Callout
add_callout(s5,
            "Mississippi leads 2 of 3 categories — lower incomes, fewer banks, highest unmet credit demand.",
            Inches(6.05), CALLOUT_V_BG, VIOLET)

# Big stat
add_text(s5, "100",
         Inches(9.5), Inches(3.0), Inches(3.0), Inches(2.0),
         size=96, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)
add_text(s5, "payday loan score\nin Mississippi\nvs. 38 national average",
         Inches(9.5), Inches(5.1), Inches(3.0), Inches(1.2),
         size=13, color=SUBTLE, align=PP_ALIGN.CENTER)

add_slide_number(s5, 5)

# ── S6: RECOMMENDATION ───────────────────────────────────────────────────────
s6 = new_slide(accent_color=GREEN)
add_tag(s6, "Recommendation")
add_h2(s6, "Go where the demand is.")

bullets = [
    ("Shift 30–40% of short-term credit budget to\nMississippi, Louisiana, Wyoming",          GREEN),
    ("Higher-intent leads in underserved markets\nwith lower advertiser competition",          SUBTLE),
    ("Lower cost-per-lead vs. saturated markets\nlike California, New York, Texas",            SUBTLE),
    ("Pipeline refreshes daily — rerun this analysis\nmonthly as demand shifts",               SUBTLE),
]
for i, (text, color) in enumerate(bullets):
    by = Inches(3.0 + i * 0.95)
    arrow = s6.shapes.add_shape(1, M, by + Inches(0.1), Inches(0.06), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GREEN
    arrow.line.fill.background()
    add_text(s6, text, M + Inches(0.2), by, Inches(10.5), Inches(0.8),
             size=17, color=color)

add_slide_number(s6, 6)

# ── S7: LIVE DEMO ─────────────────────────────────────────────────────────────
s7 = new_slide()
add_tag(s7, "Live Demo")
add_h2(s7, "Let me show you the dashboard.")
add_text(s7, "junior-data-analyst-56witkrgeceqhugwofdpzf.streamlit.app",
         M, Inches(3.3), Inches(11), Inches(0.55),
         size=16, color=INDIGO)

cols = [
    ("📈", "Descriptive tab",  "Click a state → see keyword breakdown"),
    ("🔍", "Diagnostic tab",   "Category comparison + top states"),
    ("📚", "Knowledge base",   "Industry research, synthesized"),
]
for i, (icon, title, desc) in enumerate(cols):
    cx = M + Inches(i * 3.8)
    add_text(s7, icon,  cx, Inches(4.2), Inches(3.5), Inches(0.6), size=22)
    add_text(s7, title, cx, Inches(4.9), Inches(3.5), Inches(0.5), size=15, bold=True, color=WHITE)
    add_text(s7, desc,  cx, Inches(5.5), Inches(3.5), Inches(0.6), size=13, color=SUBTLE)

add_slide_number(s7, 7)

# ── S8: CLOSE ────────────────────────────────────────────────────────────────
s8 = new_slide()
add_text(s8, "The signal exists.\nThe pipeline captures it.",
         M, Inches(1.6), Inches(11), Inches(2.2),
         size=52, bold=True, color=WHITE)
add_text(s8, "Now act on it.",
         M, Inches(3.85), Inches(11), Inches(1.0),
         size=52, bold=True, color=INDIGO)
add_text(s8, "github.com/laurendelosreyes22-commits/junior-data-analyst",
         M, Inches(5.3), Inches(11), Inches(0.55),
         size=16, color=SUBTLE)
add_slide_number(s8, 8)

# ── Save ─────────────────────────────────────────────────────────────────────
out = "docs/slides.pptx"
prs.save(out)
print(f"Saved: {out}")
